import os, sys, re, json, hashlib, io, unicodedata
from datetime import datetime
from urllib.parse import urljoin, urlparse
from collections import deque

# UTF-8 IO
try:
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")
except Exception:
    pass

from dotenv import load_dotenv, find_dotenv
load_dotenv(find_dotenv())

import requests
from requests.exceptions import SSLError
from bs4 import BeautifulSoup
from pdfminer.high_level import extract_text as pdf_extract_text
import docx
from pptx import Presentation
import openai

openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key:
    print("Erro: OPENAI_API_KEY não definida.")
    sys.exit(1)

SCRIPT_DIR  = os.path.dirname(os.path.abspath(__file__))
CHATBOT_DIR = os.path.abspath(os.path.join(SCRIPT_DIR, os.pardir, os.pardir))
OUT_BASE    = os.path.join(CHATBOT_DIR, "embeddings", "sites")

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) ChatBotCrawler/1.2"
}
TIMEOUT = 25
MAX_BIN_BYTES = 25 * 1024 * 1024  # 25 MB p/ PDF/DOCX/PPTX

# Evitar caminhos inúteis
BLOCKLIST_ROOTS = {
    "/wp-content/", "/wp-includes/", "/wp-json/", "/xmlrpc.php",
    "/feed/", "/tag/", "/author/", "/search/", "/?s="
}

# Limites anti-explosão
TOTAL_EXTERNAL_HOSTS_LIMIT = 8
PER_HOST_PAGE_LIMIT = 150  # máx. páginas processadas por host

# --------------------------
# Sanitização de nomes (sem acentos)
# --------------------------
_WIN_RESERVED = {
    "CON","PRN","AUX","NUL",
    *(f"COM{i}" for i in range(1,10)),
    *(f"LPT{i}" for i in range(1,10)),
}

def _strip_diacritics(s: str) -> str:
    # remove acentos/diacríticos preservando letras base (ç->c, ã->a etc.)
    norm = unicodedata.normalize("NFKD", s)
    return "".join(ch for ch in norm if not unicodedata.combining(ch))

def slugify(s: str) -> str:
    """
    Converte para ASCII legível:
      - remove acentos
      - troca espaços/símbolos por _
      - evita nomes reservados do Windows
    """
    s = (s or "").strip()
    if not s:
        return "Pagina"

    # alguns símbolos comuns
    s = s.replace("&", " e ")
    s = s.translate({ord("º"): "o", ord("ª"): "a", ord("°"): ""})

    # remove acentos
    s = _strip_diacritics(s)

    # mantém apenas A-Z a-z 0-9 . _ -
    s = re.sub(r"[^A-Za-z0-9._-]+", "_", s)
    s = re.sub(r"[_-]{2,}", "_", s).strip("._- ")

    if not s:
        s = "Pagina"

    head = s.split(".", 1)[0]
    if head.upper() in _WIN_RESERVED or s.startswith("."):
        s = "_" + s

    return s

def title_from_url(url: str, default="Documento") -> str:
    name = os.path.basename(urlparse(url).path) or default
    name = os.path.splitext(name)[0]
    return name or default

def make_uid_from_url(url: str) -> str:
    base = "url:" + (url or "").strip().lower()
    return hashlib.sha1(base.encode("utf-8")).hexdigest()

def clean_text_html(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    for sel in ["script", "style", "noscript"]:
        for tag in soup.select(sel):
            tag.decompose()
    for sel in ["nav", "footer", "aside"]:
        for tag in soup.select(sel):
            tag.decompose()
    text = soup.get_text(separator=" ", strip=True)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def get_title_html(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    title = soup.title.get_text(strip=True) if soup.title else ""
    if not title:
        h = soup.find(['h1','h2'])
        title = h.get_text(strip=True) if h else ""
    return title or "Página"

def canonical(u: str) -> str:
    p = urlparse(u)
    fragless = p._replace(fragment="")
    return fragless.geturl().strip()

def path_root(path: str) -> str:
    segs = [s for s in (path or "/").split("/") if s]
    return "/" + segs[0] + "/" if segs else "/"

def fetch(url: str, verify_ssl: bool = True):
    try:
        r = requests.get(url, headers=HEADERS, timeout=TIMEOUT, verify=verify_ssl)
        r.raise_for_status()
        ctype = (r.headers.get("Content-Type") or "").lower()
        return r.text, r.content, ctype, True
    except SSLError as e:
        if verify_ssl:
            print(f"[WARN SSL] {url} -> {e} (tentando sem verificação)")
            # retry 1x sem verificação
            r = requests.get(url, headers=HEADERS, timeout=TIMEOUT, verify=False)
            r.raise_for_status()
            ctype = (r.headers.get("Content-Type") or "").lower()
            return r.text, r.content, ctype, False
        raise
    except Exception:
        raise

def is_xml_like(url: str, ctype: str) -> bool:
    return ("xml" in (ctype or "")) or url.lower().endswith(".xml")

def is_pdf_like(url: str, ctype: str) -> bool:
    return ("pdf" in (ctype or "")) or url.lower().endswith(".pdf")

def is_docx_like(url: str, ctype: str) -> bool:
    return ("application/vnd.openxmlformats-officedocument.wordprocessingml.document" in (ctype or "")
            or url.lower().endswith(".docx"))

def is_pptx_like(url: str, ctype: str) -> bool:
    return ("application/vnd.openxmlformats-officedocument.presentationml.presentation" in (ctype or "")
            or url.lower().endswith(".pptx"))

def parse_sitemap(xml_text: str):
    """Extrai <loc> de urlset/sitemapindex usando parser XML; fallback regex."""
    try:
        soup = BeautifulSoup(xml_text, "xml")
        locs = [loc.get_text(strip=True) for loc in soup.find_all("loc")]
        return [u for u in locs if u and u.startswith("http")]
    except Exception:
        return re.findall(r'<loc>(.*?)</loc>', xml_text, flags=re.I)

def extract_docx_text(bin_body: bytes) -> str:
    doc = docx.Document(io.BytesIO(bin_body))
    parts = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t: parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join((cell.text or "").strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text.strip())
    txt = " ".join(parts)
    return re.sub(r"\s+"," ", txt).strip()

def extract_pptx_text(bin_body: bytes) -> str:
    prs = Presentation(io.BytesIO(bin_body))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t: parts.append(t)
    txt = " ".join(parts)
    return re.sub(r"\s+"," ", txt).strip()

def embed(text: str):
    if len(text) > 8190:  # corte conservador
        text = text[:8190]
    resp = openai.embeddings.create(
        input=text,
        model="text-embedding-ada-002"
    )
    return resp.data[0].embedding

def save_json(out_dir: str, title: str, url: str, content: str, embedding):
    os.makedirs(out_dir, exist_ok=True)
    uid = make_uid_from_url(url)
    base = slugify(title)[:120] or "Pagina"
    fn = os.path.join(out_dir, f"{base}.json")
    if os.path.exists(fn):
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        fn = os.path.join(out_dir, f"{base}_{ts}.json")
    payload = {
        "uid": uid,
        "titulo": title,
        "conteudo": content,
        "embedding": embedding,
        "link": url
    }
    with open(fn, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False)
    print(f"[JSON] {fn}")

def main():
    if len(sys.argv) < 2:
        print("Uso: crawl_site.py <url_inicial> [max_pages] [same_domain(0/1)] [use_sitemap(0/1)]")
        sys.exit(2)

    start_url   = sys.argv[1].strip()
    max_pages   = int(sys.argv[2]) if len(sys.argv) >= 3 else 10
    same_domain = int(sys.argv[3]) if len(sys.argv) >= 4 else 1
    use_sitemap = int(sys.argv[4]) if len(sys.argv) >= 5 else 0

    parsed      = urlparse(start_url)
    start_host  = parsed.netloc.lower() or "site"
    start_dir   = os.path.dirname(parsed.path or "/")
    if not start_dir.endswith("/"):
        start_dir += "/"
    start_root  = path_root(parsed.path)

    run_id = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = os.path.join(OUT_BASE, slugify(start_host), run_id)
    print(f"Coletando a partir de: {start_url}")
    print(f"Max páginas: {max_pages} | Mesmo domínio: {same_domain} | Usar sitemap: {use_sitemap}")
    print(f"Saída: {out_dir}")
    print(f"[ESCOPO] Diretório base: {start_dir}")

    # Allowlist dinâmico de HOSTS: só a partir da primeira página HTML
    external_hosts_allow = set()
    per_host_processed = {}  # host -> páginas processadas (com embed OK)

    def host_allowed(net: str) -> bool:
        net = (net or "").lower()
        if same_domain:
            return net == start_host
        # cross-domain: permite host inicial + hosts da primeira página
        return (net == start_host) or (net in external_hosts_allow)

    def can_process_host(net: str) -> bool:
        c = per_host_processed.get(net, 0)
        return c < PER_HOST_PAGE_LIMIT

    def mark_processed_host(net: str):
        per_host_processed[net] = per_host_processed.get(net, 0) + 1

    def in_scope(u: str) -> bool:
        up = urlparse(u)
        if not host_allowed(up.netloc.lower()):
            return False
        # se mesmo domínio, limita por diretório base e roots dinâmicos simples
        if up.netloc.lower() == start_host:
            # dentro do diretório inicial
            if (up.path or "/").startswith(start_dir):
                return True
            # ou dentro do mesmo root do início (ex.: /portal/, /categorias/)
            if path_root(up.path or "/") == start_root and path_root(up.path or "/") not in BLOCKLIST_ROOTS:
                return True
        # para hosts externos: não restringe por diretório, apenas pelo limite por host
        return True

    to_visit = deque()
    visited  = set()

    # Adiciona URLs do sitemap do host inicial (se houver)
    if use_sitemap:
        sm_url = f"{parsed.scheme}://{start_host}/sitemap.xml"
        try:
            sm_text, _, _, _ = fetch(sm_url)
            for u in parse_sitemap(sm_text):
                cu = canonical(u)
                if urlparse(cu).netloc.lower() == start_host and in_scope(cu):
                    to_visit.append(cu)
        except Exception as e:
            print(f"[WARN] sitemap falhou: {e}")

    if not to_visit:
        to_visit.append(canonical(start_url))

    count = 0
    first_page_hosts_locked = False  # trava allowlist depois da 1ª página

    while to_visit and count < max_pages:
        url = canonical(to_visit.popleft())
        if url in visited:
            continue
        visited.add(url)
        if not in_scope(url):
            continue

        netloc = urlparse(url).netloc.lower()
        if not can_process_host(netloc):
            print(f"[SKIP] limite por host atingido ({PER_HOST_PAGE_LIMIT}): {netloc}")
            continue

        try:
            text_body, bin_body, ctype, verified = fetch(url)
        except Exception as e:
            print(f"[ERRO] {url} -> {e}")
            continue

        # Detecta hosts externos na PRIMEIRA página HTML do host inicial
        if not first_page_hosts_locked and netloc == start_host and "html" in (ctype or ""):
            soup = BeautifulSoup(text_body, "html.parser")
            found_hosts = set()
            for a in soup.find_all("a", href=True):
                href = canonical(urljoin(url, a["href"]))
                up = urlparse(href)
                h = up.netloc.lower()
                if not h or h == start_host:
                    continue
                if h.startswith("fonts.googleapis.") or h.startswith("www.youtube.") or h.startswith("twitter.") or h.startswith("facebook."):
                    continue
                found_hosts.add(h)
            if found_hosts:
                # limita quantos hosts externos permitimos
                hosts_sorted = sorted(found_hosts)
                for h in hosts_sorted:
                    if len(external_hosts_allow) >= TOTAL_EXTERNAL_HOSTS_LIMIT:
                        break
                    external_hosts_allow.add(h)
                print("[ESCOPO] Hosts externos permitidos:", ", ".join(sorted(external_hosts_allow)) if external_hosts_allow else "(nenhum)")
            first_page_hosts_locked = True  # trava: só 1ª página define os hosts externos

        # BINÁRIOS
        if is_pdf_like(url, ctype) or is_docx_like(url, ctype) or is_pptx_like(url, ctype):
            if len(bin_body or b"") > MAX_BIN_BYTES:
                print(f"[SKIP] Arquivo grande (> {MAX_BIN_BYTES//(1024*1024)}MB): {url}")
                continue
            try:
                if is_pdf_like(url, ctype):
                    text = pdf_extract_text(io.BytesIO(bin_body)) or ""
                    title = title_from_url(url, "DocumentoPDF")
                elif is_docx_like(url, ctype):
                    text = extract_docx_text(bin_body) or ""
                    title = title_from_url(url, "DocumentoDOCX")
                else:
                    text = extract_pptx_text(bin_body) or ""
                    title = title_from_url(url, "ApresentacaoPPTX")
            except Exception as e:
                print(f"[ERRO BINARIO] {url} -> {e}")
                continue

            text = re.sub(r"\s+"," ", text.strip())
            if len(text) < 80:
                print(f"[SKIP] Arquivo com pouco texto: {url}")
                continue

            try:
                emb = embed(text)
                save_json(out_dir, title, url, text, emb)
                mark_processed_host(netloc)
                count += 1
            except Exception as e:
                print(f"[ERRO EMBED BINARIO] {url} -> {e}")
            # não segue links em binários
            continue

        # XML: só descobre URLs (apenas se do host inicial)
        if is_xml_like(url, ctype):
            if urlparse(url).netloc.lower() == start_host:
                try:
                    for u in parse_sitemap(text_body):
                        cu = canonical(u)
                        if cu not in visited and in_scope(cu):
                            to_visit.append(cu)
                except Exception:
                    pass
            continue

        # HTML
        soup = BeautifulSoup(text_body, "html.parser")
        text = clean_text_html(text_body)

        if len(text) >= 80:
            title = get_title_html(text_body)
            try:
                emb = embed(text)
                save_json(out_dir, title, url, text, emb)
                mark_processed_host(netloc)
                count += 1
            except Exception as e:
                print(f"[ERRO EMBED] {url} -> {e}")
        else:
            print(f"[SKIP] muito curto: {url}")

        # Enfileira próximos links
        try:
            for a in soup.find_all("a", href=True):
                href = urljoin(url, a["href"])
                if href.startswith(("mailto:", "javascript:")):
                    continue
                cu = canonical(href)
                up = urlparse(cu)
                # só adiciona se dentro de escopo e ainda com orçamento por host
                if cu not in visited and in_scope(cu) and can_process_host(up.netloc.lower()):
                    to_visit.append(cu)
        except Exception:
            pass

    print(f"Finalizado. Páginas processadas: {count}")
    print(f"DIR_SAIDA: {out_dir}")

if __name__ == "__main__":
    main()
